from flask import Flask, render_template, request
import requests
import PyPDF2
from pptx import Presentation
import youtube_dl

app = Flask(__name__)

def summarize_text(text):
    # Replace with your Hugging Face API URL and token
    url = "https://api-inference.huggingface.co/models/facebook/bart-large-cnn"
    headers = {"Authorization": "Bearer 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"}
    response = requests.post(url, headers=headers, json={"inputs": text})
    
    if response.status_code == 200:
        return response.json()[0]['summary_text']
    else:
        return "Error summarizing text."

def summarize_pdf(pdf_file):
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text.strip()

def summarize_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def summarize_youtube(youtube_link):
    ydl_opts = {
        'format': 'best',
        'noplaylist': True,
        'quiet': True,
        'forceurl': True,
        'postprocessors': [{
            'key': 'FFmpegExtractAudio',
            'preferredcodec': 'mp3',
            'preferredquality': '192',
        }],
    }

    with youtube_dl.YoutubeDL(ydl_opts) as ydl:
        video_info = ydl.extract_info(youtube_link, download=False)
        description = video_info.get('description', '')
        return description

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/summarize', methods=['POST'])
def summarize():
    file_type = request.form['file_type']
    simplified_text = ""
    
    if file_type == 'pdf':
        pdf_file = request.files['file']
        text = summarize_pdf(pdf_file)
        simplified_text = summarize_text(text)
        
    elif file_type == 'ppt':
        ppt_file = request.files['file']
        text = summarize_ppt(ppt_file)
        simplified_text = summarize_text(text)

    elif file_type == 'youtube':
        youtube_link = request.form['youtube_link']
        text = summarize_youtube(youtube_link)
        simplified_text = summarize_text(text)
    
    return render_template('index.html', simplified_text=simplified_text)

if __name__ == '__main__':
    app.run(debug=True)
